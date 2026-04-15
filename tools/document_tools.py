"""
High-level local document and workspace tools.
"""

from __future__ import annotations

import csv
import json
import os
import re
from html.parser import HTMLParser
from pathlib import Path
from typing import Any


_TEXT_EXTENSIONS = {
    ".md",
    ".markdown",
    ".txt",
    ".py",
    ".js",
    ".jsx",
    ".ts",
    ".tsx",
    ".json",
    ".csv",
    ".html",
    ".htm",
    ".css",
    ".yml",
    ".yaml",
    ".toml",
    ".ini",
    ".cfg",
    ".xml",
    ".log",
}
_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".bmp", ".gif", ".webp", ".tif", ".tiff"}
_BINARY_EXTENSIONS = {
    ".exe",
    ".dll",
    ".so",
    ".bin",
    ".zip",
    ".tar",
    ".gz",
    ".7z",
    ".rar",
    ".mp3",
    ".wav",
    ".mp4",
    ".mov",
    ".avi",
    ".pdf",
    ".docx",
    ".xlsx",
    ".pptx",
    *_IMAGE_EXTENSIONS,
}
_MANIFEST_NAMES = {
    "package.json",
    "package-lock.json",
    "requirements.txt",
    "pyproject.toml",
    "poetry.lock",
    "Cargo.toml",
    "go.mod",
    "pom.xml",
    "build.gradle",
    "Dockerfile",
    "docker-compose.yml",
    "docker-compose.yaml",
    "Makefile",
    "README.md",
}
_ENTRYPOINT_NAMES = {
    "main.py",
    "app.py",
    "manage.py",
    "server.py",
    "index.html",
    "index.tsx",
    "index.ts",
    "main.tsx",
    "main.ts",
    "server.js",
    "app.js",
}


class _HTMLTextExtractor(HTMLParser):
    def __init__(self):
        super().__init__()
        self._parts: list[str] = []

    def handle_data(self, data: str) -> None:
        text = data.strip()
        if text:
            self._parts.append(text)

    def get_text(self) -> str:
        return "\n".join(self._parts)


def _normalize_text(value: Any) -> str:
    return re.sub(r"\s+", " ", str(value or "").strip())


def _trim_text(value: Any, limit: int = 1200) -> str:
    normalized = _normalize_text(value)
    if len(normalized) <= limit:
        return normalized
    return normalized[: max(0, limit - 3)].rstrip() + "..."


def _safe_relpath(path: Path, root: Path) -> str:
    try:
        return str(path.relative_to(root))
    except ValueError:
        return str(path)


def _dedupe_keep_order(items: list[str]) -> list[str]:
    seen: set[str] = set()
    output: list[str] = []
    for item in items:
        normalized = str(item or "").strip()
        if not normalized or normalized in seen:
            continue
        seen.add(normalized)
        output.append(normalized)
    return output


def _is_hidden_path(path: Path) -> bool:
    return path.name.startswith(".")


def _looks_binary_path(path: Path) -> bool:
    return path.suffix.lower() in _BINARY_EXTENSIONS


def build_workspace_analysis_payload(root: Path, *, depth: int = 3, include_hidden: bool = False, focus: str = "") -> dict[str, Any]:
    if not root.exists():
        return {"tool": "analyze_workspace", "path": str(root), "status": "missing"}
    if not root.is_dir():
        return {"tool": "analyze_workspace", "path": str(root), "status": "not_directory"}

    safe_depth = max(1, min(int(depth or 3), 6))
    focus_lower = str(focus or "").strip().lower()
    tree_preview: list[str] = []
    extension_counts: dict[str, int] = {}
    manifest_files: list[str] = []
    entry_clues: list[str] = []
    large_files: list[dict[str, Any]] = []
    binary_files: list[str] = []
    focus_hits: list[str] = []
    total_files = 0
    total_dirs = 0

    def walk(current: Path, current_depth: int) -> None:
        nonlocal total_files, total_dirs
        try:
            children = sorted(current.iterdir(), key=lambda item: (item.is_file(), item.name.lower()))
        except Exception:
            return

        for child in children:
            if not include_hidden and _is_hidden_path(child):
                continue

            rel = _safe_relpath(child, root)
            tree_preview.append(f"{'  ' * current_depth}{child.name}/" if child.is_dir() else f"{'  ' * current_depth}{child.name}")
            if child.is_dir():
                total_dirs += 1
                if current_depth + 1 < safe_depth:
                    walk(child, current_depth + 1)
                continue

            total_files += 1
            suffix = child.suffix.lower() or "[no_ext]"
            extension_counts[suffix] = extension_counts.get(suffix, 0) + 1

            if child.name in _MANIFEST_NAMES:
                manifest_files.append(rel)
            if child.name in _ENTRYPOINT_NAMES:
                entry_clues.append(rel)
            if child.stat().st_size >= 5_000_000:
                large_files.append({"path": rel, "size_bytes": child.stat().st_size})
            if _looks_binary_path(child):
                binary_files.append(rel)
            if focus_lower and focus_lower in rel.lower():
                focus_hits.append(rel)

    walk(root, 0)

    extension_summary = [
        {"extension": extension, "count": count}
        for extension, count in sorted(extension_counts.items(), key=lambda item: (-item[1], item[0]))[:12]
    ]

    return {
        "tool": "analyze_workspace",
        "path": str(root),
        "status": "ok",
        "focus": _normalize_text(focus),
        "summary": {
            "directory_count": total_dirs,
            "file_count": total_files,
            "extension_counts": extension_summary,
            "manifest_files": _dedupe_keep_order(manifest_files)[:20],
            "entry_clues": _dedupe_keep_order(entry_clues)[:20],
            "large_files": sorted(large_files, key=lambda item: item["size_bytes"], reverse=True)[:10],
            "binary_files": _dedupe_keep_order(binary_files)[:20],
            "focus_hits": _dedupe_keep_order(focus_hits)[:20],
            "tree_preview": tree_preview[:80],
        },
        "answer_style": "Use the workspace summary to explain the directory structure, likely entrypoints, manifests, and potential hotspots.",
    }


class DocumentTools:
    def __init__(self, mode_manager, agent_dispatcher=None, allow_local_fallback: bool = True):
        self._mode_manager = mode_manager
        self._agent_dispatcher = agent_dispatcher
        self._allow_local_fallback = bool(allow_local_fallback)

    def set_agent_dispatcher(self, dispatcher) -> None:
        self._agent_dispatcher = dispatcher

    def set_capability_dispatcher(self, dispatcher) -> None:
        self.set_agent_dispatcher(dispatcher)

    def set_local_fallback_enabled(self, enabled: bool) -> None:
        self._allow_local_fallback = bool(enabled)

    async def _dispatch_capability(
        self,
        *,
        capability_suffix: str,
        arguments: dict[str, Any],
        session_id: str,
        title: str,
        operation_type: str,
    ) -> dict[str, Any]:
        dispatcher = self._agent_dispatcher
        if dispatcher is None:
            raise RuntimeError("Capability dispatcher is not configured")
        dispatch = getattr(dispatcher, "dispatch_agent_capability", None)
        if not callable(dispatch):
            dispatch = getattr(dispatcher, "dispatch_local_capability", None)
        if not callable(dispatch):
            raise RuntimeError("Capability dispatcher does not support capability dispatch")
        return await dispatch(
            capability_suffix=capability_suffix,
            arguments=arguments,
            session_id=session_id,
            title=title,
            operation_type=operation_type,
        )

    def _ensure_local_capability_available(
        self,
        *,
        capability_suffix: str,
        session_id: str = "",
        path: str = "",
    ) -> None:
        if self._agent_dispatcher is not None or self._allow_local_fallback:
            return
        error = RuntimeError("Core local fallback is disabled")
        error.tool_error_code = "local_agent_required"
        error.tool_error_message = "当前 Core 不再直接执行本地文档能力，请连接 Desktop Agent 后重试。"
        error.tool_error_details = {
            "capability_suffix": capability_suffix,
            "session_id": str(session_id or ""),
            "path": str(path or ""),
        }
        error.tool_error_retryable = False
        raise error

    def _parser_config(self) -> dict[str, Any]:
        return self._mode_manager.get_document_parser_config()

    def _resolve_path(self, path_value: str) -> Path:
        raw = str(path_value or "").strip()
        if not raw:
            raise ValueError("path is required")
        return Path(raw).expanduser().resolve()

    def _is_hidden(self, path: Path) -> bool:
        return _is_hidden_path(path)

    def _looks_binary(self, path: Path) -> bool:
        return _looks_binary_path(path)

    def _chunk_text(self, text: str, chunking: str, *, max_chunks: int) -> list[dict[str, Any]]:
        normalized = str(text or "").strip()
        if not normalized:
            return []

        if chunking == "none":
            return [{"index": 1, "label": "full", "content": _trim_text(normalized, 2400)}]

        if chunking in {"paragraph", "auto"}:
            parts = [part.strip() for part in re.split(r"\n\s*\n", normalized) if part.strip()]
        else:
            parts = [line.strip() for line in normalized.splitlines() if line.strip()]

        if not parts:
            parts = [normalized]

        chunks: list[dict[str, Any]] = []
        for index, part in enumerate(parts[:max_chunks], start=1):
            chunks.append(
                {
                    "index": index,
                    "label": f"chunk_{index}",
                    "content": _trim_text(part, 1800),
                }
            )
        return chunks

    def _read_text_file(self, path: Path) -> str:
        for encoding in ("utf-8", "utf-8-sig", "gb18030", "latin-1"):
            try:
                return path.read_text(encoding=encoding)
            except UnicodeDecodeError:
                continue
        return path.read_text(encoding="utf-8", errors="replace")

    def _read_content_from_text(self, path: Path, raw_text: str) -> tuple[str, str]:
        suffix = path.suffix.lower()
        if suffix == ".json":
            payload = json.loads(raw_text)
            return "json", json.dumps(payload, ensure_ascii=False, indent=2)
        if suffix == ".csv":
            rows: list[list[str]] = []
            reader = csv.reader(str(raw_text or "").splitlines())
            for index, row in enumerate(reader):
                rows.append([str(item) for item in row])
                if index >= 20:
                    break
            return "csv", "\n".join(", ".join(row) for row in rows)
        if suffix in {".html", ".htm"}:
            extractor = _HTMLTextExtractor()
            extractor.feed(str(raw_text or ""))
            return "html", extractor.get_text()
        return "text", str(raw_text or "")

    def _read_json_file(self, path: Path) -> str:
        payload = json.loads(self._read_text_file(path))
        return json.dumps(payload, ensure_ascii=False, indent=2)

    def _read_csv_file(self, path: Path) -> str:
        rows: list[list[str]] = []
        with path.open("r", encoding="utf-8", errors="replace", newline="") as handle:
            reader = csv.reader(handle)
            for index, row in enumerate(reader):
                rows.append([str(item) for item in row])
                if index >= 20:
                    break
        return "\n".join(", ".join(row) for row in rows)

    def _read_html_file(self, path: Path) -> str:
        extractor = _HTMLTextExtractor()
        extractor.feed(self._read_text_file(path))
        return extractor.get_text()

    def _read_pdf_file(self, path: Path) -> str:
        try:
            import fitz  # type: ignore

            doc = fitz.open(path)
            try:
                parts = [doc.load_page(index).get_text("text") for index in range(min(doc.page_count, 20))]
            finally:
                doc.close()
            return "\n\n".join(parts)
        except ImportError:
            try:
                import pdfplumber  # type: ignore

                with pdfplumber.open(path) as pdf:
                    return "\n\n".join(
                        (page.extract_text() or "")
                        for page in pdf.pages[:20]
                    )
            except ImportError as exc:
                raise RuntimeError("PDF parser unavailable. Install PyMuPDF or pdfplumber.") from exc

    def _read_docx_file(self, path: Path) -> str:
        try:
            from docx import Document  # type: ignore
        except ImportError as exc:
            raise RuntimeError("DOCX parser unavailable. Install python-docx.") from exc

        document = Document(path)
        parts = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        return "\n".join(parts)

    def _read_xlsx_file(self, path: Path) -> str:
        try:
            from openpyxl import load_workbook  # type: ignore
        except ImportError as exc:
            raise RuntimeError("XLSX parser unavailable. Install openpyxl.") from exc

        workbook = load_workbook(path, read_only=True, data_only=True)
        try:
            parts: list[str] = []
            for sheet in workbook.worksheets[:5]:
                parts.append(f"[Sheet] {sheet.title}")
                for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                    values = [str(item) for item in row if item not in (None, "")]
                    if values:
                        parts.append(", ".join(values))
                    if row_index >= 20:
                        break
            return "\n".join(parts)
        finally:
            workbook.close()

    def _read_pptx_file(self, path: Path) -> str:
        try:
            from pptx import Presentation  # type: ignore
        except ImportError as exc:
            raise RuntimeError("PPTX parser unavailable. Install python-pptx.") from exc

        presentation = Presentation(path)
        parts: list[str] = []
        for slide_index, slide in enumerate(presentation.slides[:20], start=1):
            parts.append(f"[Slide] {slide_index}")
            for shape in slide.shapes:
                if hasattr(shape, "text") and str(shape.text).strip():
                    parts.append(str(shape.text).strip())
        return "\n".join(parts)

    def _read_image_ocr(self, path: Path) -> str:
        if not self._parser_config().get("enable_ocr", True):
            raise RuntimeError("OCR parser disabled by configuration.")
        try:
            from PIL import Image  # type: ignore
            import pytesseract  # type: ignore
        except ImportError as exc:
            raise RuntimeError("OCR parser unavailable. Install Pillow and pytesseract.") from exc

        image = Image.open(path)
        try:
            return pytesseract.image_to_string(image)
        finally:
            image.close()

    def _read_document_content(self, path: Path) -> tuple[str, str]:
        suffix = path.suffix.lower()
        if suffix == ".json":
            return "json", self._read_json_file(path)
        if suffix == ".csv":
            return "csv", self._read_csv_file(path)
        if suffix in {".html", ".htm"}:
            return "html", self._read_html_file(path)
        if suffix == ".pdf":
            return "pdf", self._read_pdf_file(path)
        if suffix == ".docx":
            return "docx", self._read_docx_file(path)
        if suffix == ".xlsx":
            return "xlsx", self._read_xlsx_file(path)
        if suffix == ".pptx":
            return "pptx", self._read_pptx_file(path)
        if suffix in _IMAGE_EXTENSIONS:
            return "image_ocr", self._read_image_ocr(path)
        return "text", self._read_text_file(path)

    def _collect_document(
        self,
        path: Path,
        *,
        goal: str,
        chunking: str,
    ) -> dict[str, Any]:
        config = self._parser_config()
        max_bytes = int(config.get("max_file_bytes", 2_000_000) or 2_000_000)
        max_total_chars = int(config.get("max_total_chars", 24_000) or 24_000)
        max_chunks = int(config.get("max_chunks_per_document", 12) or 12)

        if not path.exists():
            return {"path": str(path), "status": "missing"}
        if path.is_dir():
            return {
                "path": str(path),
                "status": "directory",
                "warning": "Use analyze_workspace for directories.",
            }

        file_size = path.stat().st_size
        if file_size > max_bytes:
            return {
                "path": str(path),
                "status": "too_large",
                "size_bytes": file_size,
                "warning": f"File exceeds parser limit of {max_bytes} bytes.",
            }

        try:
            doc_type, raw_content = self._read_document_content(path)
        except Exception as exc:
            return {
                "path": str(path),
                "status": "unreadable",
                "size_bytes": file_size,
                "warning": str(exc),
            }

        content = str(raw_content or "").strip()
        trimmed = content[:max_total_chars]
        truncated = len(content) > len(trimmed)
        chunks = self._chunk_text(trimmed, chunking, max_chunks=max_chunks)

        return {
            "path": str(path),
            "name": path.name,
            "type": doc_type,
            "status": "ok",
            "size_bytes": file_size,
            "goal": _normalize_text(goal),
            "truncated": truncated,
            "char_count": len(content),
            "content_excerpt": _trim_text(trimmed, 2000),
            "chunks": chunks,
        }

    async def _collect_document_via_agent(
        self,
        path: Path,
        *,
        goal: str,
        chunking: str,
        session_id: str,
    ) -> dict[str, Any]:
        config = self._parser_config()
        max_total_chars = int(config.get("max_total_chars", 24_000) or 24_000)
        max_chunks = int(config.get("max_chunks_per_document", 12) or 12)
        try:
            result = await self._dispatch_capability(
                capability_suffix="file.read",
                arguments={"path": str(path), "max_chars": max_total_chars},
                session_id=session_id,
                title=f"Read Local File: {path.name}",
                operation_type="tool.read_local_documents",
            )
            doc_type, raw_content = self._read_content_from_text(path, str(result.get("content") or ""))
        except Exception as exc:
            return {
                "path": str(path),
                "status": "unreadable",
                "warning": str(exc),
            }

        content = str(raw_content or "").strip()
        trimmed = content[:max_total_chars]
        truncated = len(content) > len(trimmed) or bool(result.get("truncated"))
        chunks = self._chunk_text(trimmed, chunking, max_chunks=max_chunks)
        return {
            "path": str(path),
            "name": path.name,
            "type": doc_type,
            "status": "ok",
            "size_bytes": int(result.get("size_bytes") or 0),
            "goal": _normalize_text(goal),
            "truncated": truncated,
            "char_count": len(content),
            "content_excerpt": _trim_text(trimmed, 2000),
            "chunks": chunks,
        }

    async def analyze_workspace(
        self,
        path: str,
        depth: int = 3,
        include_hidden: bool = False,
        focus: str = "",
        session_id: str = "",
        source=None,
        route_context: dict[str, Any] | None = None,
        activity_callback=None,
    ) -> str:
        del source, route_context, activity_callback
        root = self._resolve_path(path)
        self._ensure_local_capability_available(
            capability_suffix="workspace.analyze",
            session_id=session_id,
            path=str(root),
        )
        if self._agent_dispatcher is None:
            payload = build_workspace_analysis_payload(root, depth=depth, include_hidden=include_hidden, focus=focus)
        else:
            payload = await self._dispatch_capability(
                capability_suffix="workspace.analyze",
                arguments={
                    "path": str(root),
                    "depth": depth,
                    "include_hidden": include_hidden,
                    "focus": focus,
                },
                session_id=session_id,
                title=f"Analyze Workspace: {root.name}",
                operation_type="tool.analyze_workspace",
            )
        return json.dumps(payload, ensure_ascii=False, indent=2)

    async def read_local_documents(
        self,
        paths: list[str] | str,
        goal: str = "",
        chunking: str = "auto",
        session_id: str = "",
        source=None,
        route_context: dict[str, Any] | None = None,
        activity_callback=None,
    ) -> str:
        del source, route_context, activity_callback
        normalized_paths = [paths] if isinstance(paths, str) else list(paths or [])
        self._ensure_local_capability_available(
            capability_suffix="file.read",
            session_id=session_id,
            path=",".join(str(self._resolve_path(item)) for item in normalized_paths[:5]),
        )
        if self._agent_dispatcher is None:
            documents = [self._collect_document(self._resolve_path(item), goal=goal, chunking=chunking) for item in normalized_paths]
        else:
            documents = [
                await self._collect_document_via_agent(
                    self._resolve_path(item),
                    goal=goal,
                    chunking=chunking,
                    session_id=session_id,
                )
                for item in normalized_paths
            ]
        readable = [item for item in documents if item.get("status") == "ok"]
        payload = {
            "tool": "read_local_documents",
            "goal": _normalize_text(goal),
            "chunking": str(chunking or "auto"),
            "document_count": len(documents),
            "readable_count": len(readable),
            "documents": documents,
            "answer_style": "Use only the returned document excerpts and chunks when summarizing or quoting local files.",
        }
        return json.dumps(payload, ensure_ascii=False, indent=2)

    async def write_local_document(
        self,
        path: str,
        content: str,
        mode: str = "overwrite",
        preview: bool = True,
        confirmed: bool = False,
        session_id: str = "",
        source=None,
        route_context: dict[str, Any] | None = None,
        activity_callback=None,
    ) -> str:
        del source, route_context, activity_callback, confirmed
        target = self._resolve_path(path)
        normalized_mode = str(mode or "overwrite").strip().lower() or "overwrite"
        trusted = self._mode_manager.is_trusted_write_path(str(target))
        payload = {
            "tool": "write_local_document",
            "path": str(target),
            "mode": normalized_mode,
            "trusted_root": trusted,
            "action_risk": "local_write",
            "preview": bool(preview),
            "content_preview": _trim_text(content, 600),
            "byte_length": len(str(content or "").encode("utf-8")),
        }

        if preview:
            payload["status"] = "preview"
            payload["answer_style"] = "Ask for confirmation or call again with preview=false to perform the write."
            return json.dumps(payload, ensure_ascii=False, indent=2)

        if not trusted:
            payload["status"] = "blocked_untrusted_path"
            payload["trusted_write_roots"] = self._mode_manager.get_trusted_write_roots()
            return json.dumps(payload, ensure_ascii=False, indent=2)

        self._ensure_local_capability_available(
            capability_suffix="file.write",
            session_id=session_id,
            path=str(target),
        )
        if self._agent_dispatcher is not None:
            result = await self._dispatch_capability(
                capability_suffix="file.write",
                arguments={"path": str(target), "content": str(content or ""), "mode": normalized_mode},
                session_id=session_id,
                title=f"Write Local File: {target.name}",
                operation_type="tool.write_local_document",
            )
            payload["status"] = "written"
            payload["bytes_written"] = result.get("bytes_written", 0)
            return json.dumps(payload, ensure_ascii=False, indent=2)

        target.parent.mkdir(parents=True, exist_ok=True)
        if normalized_mode == "append":
            with target.open("a", encoding="utf-8") as handle:
                handle.write(str(content or ""))
        elif normalized_mode == "create_if_missing":
            if target.exists():
                payload["status"] = "already_exists"
                return json.dumps(payload, ensure_ascii=False, indent=2)
            target.write_text(str(content or ""), encoding="utf-8")
        else:
            target.write_text(str(content or ""), encoding="utf-8")

        payload["status"] = "written"
        return json.dumps(payload, ensure_ascii=False, indent=2)

    def _replace_markdown_section(self, text: str, selector: str, replacement_content: str) -> tuple[str, bool]:
        lines = text.splitlines()
        selector_lower = str(selector or "").strip().lower()
        if not selector_lower:
            return replacement_content, True

        heading_pattern = re.compile(r"^(#+)\s*(.+?)\s*$")
        start_index = -1
        end_index = len(lines)
        level = 1
        for index, line in enumerate(lines):
            match = heading_pattern.match(line)
            if not match:
                continue
            heading_text = match.group(2).strip().lower()
            if selector_lower in heading_text:
                start_index = index
                level = len(match.group(1))
                break

        if start_index == -1:
            return text, False

        for index in range(start_index + 1, len(lines)):
            match = heading_pattern.match(lines[index])
            if match and len(match.group(1)) <= level:
                end_index = index
                break

        replacement_lines = [lines[start_index], str(replacement_content or "").strip()]
        new_lines = lines[:start_index] + replacement_lines + lines[end_index:]
        return "\n".join(new_lines).strip() + "\n", True

    async def rewrite_local_document(
        self,
        path: str,
        instructions: str,
        section_selector: str = "",
        preview: bool = True,
        replacement_content: str = "",
        confirmed: bool = False,
        session_id: str = "",
        source=None,
        route_context: dict[str, Any] | None = None,
        activity_callback=None,
    ) -> str:
        del source, route_context, activity_callback, confirmed
        target = self._resolve_path(path)
        self._ensure_local_capability_available(
            capability_suffix="file.read",
            session_id=session_id,
            path=str(target),
        )
        if self._agent_dispatcher is None and (not target.exists() or target.is_dir()):
            return json.dumps(
                {
                    "tool": "rewrite_local_document",
                    "path": str(target),
                    "status": "missing_or_not_file",
                    "action_risk": "read",
                },
                ensure_ascii=False,
                indent=2,
            )

        if self._agent_dispatcher is None:
            original = self._read_text_file(target)
        else:
            try:
                read_result = await self._dispatch_capability(
                    capability_suffix="file.read",
                    arguments={"path": str(target), "max_chars": 200000},
                    session_id=session_id,
                    title=f"Read For Rewrite: {target.name}",
                    operation_type="tool.rewrite_local_document.read",
                )
                original = str(read_result.get("content") or "")
            except Exception:
                return json.dumps(
                    {
                        "tool": "rewrite_local_document",
                        "path": str(target),
                        "status": "missing_or_not_file",
                        "action_risk": "read",
                    },
                    ensure_ascii=False,
                    indent=2,
                )
        payload = {
            "tool": "rewrite_local_document",
            "path": str(target),
            "section_selector": _normalize_text(section_selector),
            "instructions": _normalize_text(instructions),
            "trusted_root": self._mode_manager.is_trusted_write_path(str(target)),
        }

        if not replacement_content:
            payload["status"] = "needs_replacement_content"
            payload["action_risk"] = "read"
            payload["current_excerpt"] = _trim_text(original, 1600)
            payload["answer_style"] = (
                "Generate the revised content first, then call rewrite_local_document again with replacement_content, "
                "or call write_local_document with the final full document."
            )
            return json.dumps(payload, ensure_ascii=False, indent=2)

        if section_selector:
            rewritten, replaced = self._replace_markdown_section(original, section_selector, replacement_content)
            if not replaced:
                payload["status"] = "section_not_found"
                payload["action_risk"] = "read"
                payload["current_excerpt"] = _trim_text(original, 1600)
                return json.dumps(payload, ensure_ascii=False, indent=2)
        else:
            rewritten = str(replacement_content or "")

        payload["action_risk"] = "local_write"
        payload["preview"] = bool(preview)
        payload["proposed_excerpt"] = _trim_text(rewritten, 1600)

        if preview:
            payload["status"] = "preview"
            return json.dumps(payload, ensure_ascii=False, indent=2)

        if not self._mode_manager.is_trusted_write_path(str(target)):
            payload["status"] = "blocked_untrusted_path"
            payload["trusted_write_roots"] = self._mode_manager.get_trusted_write_roots()
            return json.dumps(payload, ensure_ascii=False, indent=2)

        if self._agent_dispatcher is not None:
            await self._dispatch_capability(
                capability_suffix="file.write",
                arguments={"path": str(target), "content": rewritten, "mode": "overwrite"},
                session_id=session_id,
                title=f"Rewrite Local File: {target.name}",
                operation_type="tool.rewrite_local_document.write",
            )
            payload["status"] = "written"
            return json.dumps(payload, ensure_ascii=False, indent=2)

        target.write_text(rewritten, encoding="utf-8")
        payload["status"] = "written"
        return json.dumps(payload, ensure_ascii=False, indent=2)

    async def compile_report(
        self,
        inputs: list[Any],
        format: str = "markdown",
        template: str = "report",
        title: str = "",
        session_id: str = "",
        source=None,
        route_context: dict[str, Any] | None = None,
        activity_callback=None,
    ) -> str:
        del session_id, source, route_context, activity_callback
        sections: list[dict[str, Any]] = []
        for index, item in enumerate(list(inputs or []), start=1):
            if isinstance(item, dict):
                section_title = _normalize_text(item.get("title") or f"Section {index}")
                section_content = _trim_text(item.get("content") or item.get("summary") or "", 2400)
                sections.append({"title": section_title, "content": section_content})
                continue

            text = str(item or "").strip()
            if not text:
                continue

            maybe_path = Path(text).expanduser()
            if maybe_path.exists() and maybe_path.is_file():
                collected = self._collect_document(maybe_path.resolve(), goal="", chunking="none")
                section_content = collected.get("content_excerpt") or collected.get("warning") or ""
                section_title = collected.get("name") or maybe_path.name
            else:
                section_title = f"Input {index}"
                section_content = _trim_text(text, 2400)
            sections.append({"title": section_title, "content": section_content})

        normalized_template = str(template or "report").strip().lower() or "report"
        normalized_format = str(format or "markdown").strip().lower() or "markdown"
        report_title = _normalize_text(title) or normalized_template.replace("_", " ").title()

        if normalized_format == "json":
            rendered = json.dumps(
                {
                    "title": report_title,
                    "template": normalized_template,
                    "sections": sections,
                },
                ensure_ascii=False,
                indent=2,
            )
        elif normalized_format == "html":
            rendered = "\n".join(
                [
                    f"<h1>{report_title}</h1>",
                    *[
                        f"<section><h2>{section['title']}</h2><p>{section['content']}</p></section>"
                        for section in sections
                    ],
                ]
            )
        else:
            lines = [f"# {report_title}", "", f"_Template: {normalized_template}_", ""]
            for section in sections:
                lines.extend([f"## {section['title']}", section["content"], ""])
            rendered = "\n".join(lines).strip() + "\n"

        payload = {
            "tool": "compile_report",
            "template": normalized_template,
            "format": normalized_format,
            "title": report_title,
            "section_count": len(sections),
            "sections": sections,
            "compiled_report": rendered,
            "answer_style": "Use the compiled report directly or write it with write_local_document if the user wants a saved file.",
        }
        return json.dumps(payload, ensure_ascii=False, indent=2)
